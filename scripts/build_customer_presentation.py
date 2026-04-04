"""
Generate customer-facing Parking Platform presentation (.pptx).
Run from repo root: pip install python-pptx && python scripts/build_customer_presentation.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

OUTPUT = Path(__file__).resolve().parent.parent / "Parking_Platform_Customer_Presentation.pptx"


def _title_only(prs: Presentation, title: str, subtitle: str | None = None) -> None:
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def _bullets(prs: Presentation, title: str, lines: list[str]) -> None:
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(lines):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)


def _blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            return layout
    return prs.slide_layouts[-1]


def _table_slide(prs: Presentation, title: str, headers: tuple[str, ...], rows: list[tuple[str, str]]) -> None:
    slide = prs.slides.add_slide(_blank_layout(prs))
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.8))
    tf = tbox.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True

    nrows = 1 + len(rows)
    ncols = len(headers)
    left, top, width, height = Inches(0.5), Inches(1.25), Inches(9), Inches(5.2)
    table = slide.shapes.add_table(nrows, ncols, left, top, width, height).table

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(14)

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _title_only(
        prs,
        "Smart parking operations platform",
        "Garage configuration, live occupancy, tickets, payments, and reporting",
    )
    _title_only(
        prs,
        "Parking API + Garage Dashboard",
        "PostgreSQL-backed REST API · Vue 3 operator UI · Production-oriented security",
    )

    _bullets(
        prs,
        "The problem we solve",
        [
            "Parking operators need accurate spot occupancy and open sessions.",
            "They need reliable entry/exit workflows and consistent pricing.",
            "They need clear payment handling and revenue visibility.",
            "Systems must scale to multiple garages and integrations.",
            "This solution delivers an operational backbone (API + data) and a modern web dashboard.",
        ],
    )

    _table_slide(
        prs,
        "What you get",
        ("Area", "Capability"),
        [
            ("Garages", "Capacity, rates (default, day/night, lost ticket), hours, subscriptions"),
            ("Spots", "Per-garage spots; rentable/active; occupancy tied to tickets"),
            ("Vehicles", "Registry (plate, type); supports entry workflows"),
            ("Tickets", "Entry/exit, states, fees, payment and operational status"),
            ("Payments", "Record payments for closed tickets; currency support (e.g. RSD)"),
            ("Media", "Ticket image upload (client resizes before upload)"),
            ("Dashboard", "Status, garage overview, ticket activity, revenue, per-garage drill-down"),
        ],
    )

    _bullets(
        prs,
        "Operator dashboard",
        [
            "Secure login with JWT; protected routes for operators.",
            "Live sync: periodic refresh plus refresh when returning to the browser tab.",
            "Garage overview and per-garage detail: spots, open tickets, activity, revenue metrics.",
            "New vehicle entry from the UI; ticket detail and payment flows.",
            "Bilingual interface: English and Serbian.",
        ],
    )

    _bullets(
        prs,
        "Technical architecture",
        [
            "Browser: Vue 3 + TypeScript + Vite + Tailwind dashboard.",
            "Backend: FastAPI REST API with OpenAPI (Swagger) documentation.",
            "Data: PostgreSQL; schema migrations via Alembic.",
            "Uploaded ticket images stored on the server and served securely.",
            "Health endpoint verifies application and database connectivity.",
        ],
    )

    _bullets(
        prs,
        "Security and integration",
        [
            "JWT authentication after login; optional API key for tools and integrations.",
            "CORS configurable for production web origins.",
            "Same API serves the dashboard, future mobile apps, or partner systems.",
        ],
    )

    _bullets(
        prs,
        "Deployment flexibility",
        [
            "Supports database triggers for fees and payment status, or API-side calculation.",
            "Useful for greenfield installs, migrations, or existing PostgreSQL rules.",
            "Backend: Uvicorn; configurable host, port, and environment variables.",
            "Frontend: static build output deployable to any CDN or reverse proxy.",
        ],
    )

    _bullets(
        prs,
        "Quality and maintainability",
        [
            "Integration tests (pytest) against the real API and database.",
            "Structured codebase: routers, services, models, and schemas.",
            "Consistent API errors and documented endpoints for integrators.",
        ],
    )

    _bullets(
        prs,
        "Summary and next steps",
        [
            "A parking garage operations platform: secure REST API, PostgreSQL, real-time dashboard.",
            "Suited to single or multi-garage operations and future integrations.",
            "Next: live demo and alignment on pricing rules, garages, and hosting.",
        ],
    )

    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    main()
