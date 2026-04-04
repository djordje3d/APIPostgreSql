"""
Generate project / application structure presentation (.pptx).

Run: python -m pip install python-pptx
     python scripts/build_structure_presentation.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

OUTPUT = (
    Path(__file__).resolve().parent.parent
    / "Parking_Platform_Project_Structure.pptx"
)


def _title_only(
    prs: Presentation,
    title: str,
    subtitle: str | None = None,
) -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def _bullets(
    prs: Presentation,
    title: str,
    lines: list[str],
    font_pt: int = 18,
) -> None:
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(lines):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(font_pt)


def _blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if layout.name == "Blank":
            return layout
    return prs.slide_layouts[-1]


def _code_block_slide(
    prs: Presentation,
    title: str,
    body_text: str,
    font_pt: int = 13,
) -> None:
    slide = prs.slides.add_slide(_blank_layout(prs))
    tbox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(9), Inches(0.65)
    )
    tf = tbox.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True

    box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.05), Inches(9), Inches(6.2)
    )
    btf = box.text_frame
    btf.clear()
    first = True
    for line in body_text.strip().split("\n"):
        p = btf.paragraphs[0] if first else btf.add_paragraph()
        first = False
        p.text = line
        p.font.size = Pt(font_pt)
        p.font.name = "Consolas"


def _table_slide(
    prs: Presentation,
    title: str,
    headers: tuple[str, ...],
    rows: list[tuple[str, str]],
) -> None:
    slide = prs.slides.add_slide(_blank_layout(prs))
    tbox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(9), Inches(0.8)
    )
    tf = tbox.text_frame
    tf.paragraphs[0].text = title
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True

    nrows, ncols = 1 + len(rows), len(headers)
    table = slide.shapes.add_table(
        nrows,
        ncols,
        Inches(0.5),
        Inches(1.2),
        Inches(9),
        Inches(5.9),
    ).table

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(13)

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    _title_only(
        prs,
        "Parking platform — project structure",
        (
            "Repository layout: backend (app), frontend (dashboard), "
            "database migrations, tests"
        ),
    )

    _code_block_slide(
        prs,
        "Repository (top level)",
        """
APIPostgreSql/
├── app/                 FastAPI backend (Python package)
├── dashboard/           Vue 3 + Vite + TypeScript SPA
├── tests/               Pytest API integration tests
├── alembic/             SQLAlchemy / Alembic migrations
├── scripts/             Dev tooling (e.g. presentation generators)
├── requirements.txt     Python dependencies
├── README.md            API setup, env vars, run & test commands
└── (env, Postman docs, etc.)
        """,
        font_pt=14,
    )

    _bullets(
        prs,
        "Application architecture (runtime)",
        [
            (
                "Browser runs the dashboard (Vue); it calls the REST API "
                "over HTTP(S)."
            ),
            (
                "FastAPI app (app.main) mounts routers, CORS, auth "
                "middleware, static /uploads."
            ),
            (
                "Routers validate input (Pydantic schemas) and call services "
                "for business logic."
            ),
            "Services use SQLAlchemy sessions to read/write PostgreSQL.",
            (
                "Ticket images: upload endpoint writes files; /uploads "
                "serves them."
            ),
        ],
        font_pt=16,
    )

    _code_block_slide(
        prs,
        "Backend package: app/",
        """
app/
├── main.py              FastAPI app, middleware, route registration, /health
├── run.py               Uvicorn entry point
├── config.py            Environment (.env), feature flags, CORS, paths
├── db.py                Engine, session factory, get_db dependency
├── models.py            SQLAlchemy ORM models (garages, spots, vehicles, …)
├── schemas.py           Pydantic request/response models
├── auth.py              API key middleware (X-API-Key or Bearer JWT)
├── auth_jwt.py          JWT create/verify, login dependencies
├── errors.py            Structured API error helper
├── error_handlers.py    Exception → HTTP responses
├── routers/             HTTP endpoints per domain (thin handlers)
├── services/            Business logic (tickets, spots, payments, pricing, …)
├── sql/                 Ad-hoc SQL snippets (analytics / maintenance)
└── migrations/          Optional one-off SQL (e.g. image column notes)
        """,
        font_pt=11,
    )

    _table_slide(
        prs,
        "app/routers — HTTP surface",
        ("Module", "Role"),
        [
            ("auth.py", "POST /auth/login → JWT"),
            ("garages.py", "Parking garage CRUD / config"),
            ("vehicle_types.py", "Vehicle type definitions & rates"),
            ("vehicles.py", "Vehicle registry"),
            ("tickets.py", "Entry, exit, listing, ticket lifecycle"),
            ("payments.py", "Payments for tickets"),
            ("spots.py", "Parking spots per garage"),
            ("dashboard.py", "Aggregated analytics (fewer round-trips)"),
            ("upload.py", "Ticket image upload"),
        ],
    )

    _table_slide(
        prs,
        "app/services — business logic",
        ("Module", "Role"),
        [
            ("tickets.py", "Ticket flow, allocation, state/fee coordination"),
            ("spots.py", "Spot listing / updates used by API"),
            ("payments.py", "Payment recording, status updates"),
            ("pricing.py", "Fee / rate calculations when API-driven"),
            ("dashboard_analytics.py", "Counts, revenue sums, outstanding"),
            ("tokens.py", "Ticket token generation"),
            ("rewrite_ticket_tokens.py", "Token migration helpers"),
        ],
    )

    _code_block_slide(
        prs,
        "Frontend: dashboard/",
        """
dashboard/
├── package.json, vite.config.ts, tailwind config
├── env.example          VITE_API_URL, optional VITE_API_KEY
├── src/
│   ├── main.ts          Bootstrap Vue + router + i18n
│   ├── App.vue          Shell layout
│   ├── router/index.ts  /, /login, /garage/:id (auth guard)
│   ├── i18n.ts, locales/   en.json, sr.json
│   ├── api/             Typed fetch wrappers (client, auth, tickets, …)
│   ├── views/           DashboardView, GarageDetailView, LoginView
│   ├── components/
│   │   ├── dashboard/   Garage tables, tickets, charts, modals
│   │   └── ui/          Modal, inputs, toast, pagination, …
│   ├── composables/     Polling, garage hooks, formatters, toast
│   ├── services/        Orchestration (e.g. create parking entry)
│   ├── utils/           Dates, cache, images, barcode helpers
│   └── constants/       Refresh intervals, shared constants
└── dist/                Production build output (npm run build)
        """,
        font_pt=10,
    )

    _bullets(
        prs,
        "Dashboard — data flow",
        [
            "api/client.ts — base URL, credentials, error handling.",
            (
                "Domain modules (tickets.ts, garages.ts, …) map to backend "
                "routes."
            ),
            (
                "Composables poll dashboard analytics and garage endpoints "
                "on a timer."
            ),
            (
                "Views compose dashboard/* components; modals handle entry, "
                "payment, ticket detail."
            ),
        ],
        font_pt=17,
    )

    _code_block_slide(
        prs,
        "tests/ and alembic/",
        """
tests/
├── conftest.py          Fixtures, DB session per test (rollback)
├── test_health.py       /health
├── test_auth.py         Login / API key behavior
├── test_garages.py
├── test_vehicle_types.py
├── test_vehicles.py
├── test_spots.py
├── test_tickets_flow.py
└── test_payments.py

alembic/
├── env.py               Migration environment (DB URL / models)
├── versions/            Revision scripts (schema history)
└── script.py.mako       Template for new revisions
        """,
        font_pt=12,
    )

    _bullets(
        prs,
        "Summary",
        [
            (
                "Clear split: app = API + domain logic; "
                "dashboard = operator UI."
            ),
            (
                "Routers stay thin; services hold rules; models/schemas "
                "define contracts."
            ),
            (
                "Alembic versions database schema; tests lock in API "
                "behavior."
            ),
            (
                "Regenerate after major moves: "
                "python scripts/build_structure_presentation.py"
            ),
        ],
        font_pt=17,
    )

    prs.save(OUTPUT)
    print(f"Wrote {OUTPUT}")


if __name__ == "__main__":
    main()
